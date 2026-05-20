import base64
import re
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote

import httpx
from fastapi import APIRouter, Depends, HTTPException, Query
from fastapi.responses import FileResponse
from loguru import logger
from pptx import Presentation
from pptx.util import Emu, Inches
from sqlmodel import Session

from app.core.auth import get_current_user_id
from app.core.image_storage import URL_PREFIX as IMAGES_URL_PREFIX, resolve_local_path
from app.models.db import get_session
from app.models.schemas import ExportRequest, ExportResponse
from app.services.export_service import ExportService
from app.services.project_service import ProjectService

router = APIRouter(tags=["export"])


@router.post("/api/projects/{project_id}/export", response_model=ExportResponse)
def export_project(
    project_id: str,
    request: ExportRequest,
    session: Session = Depends(get_session),
    user_id: int = Depends(get_current_user_id),
) -> ExportResponse:
    ProjectService(session)._get_owned_record(project_id, user_id)
    try:
        return ExportService(session).export_project(project_id, request.format, request.include_index)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@router.get("/api/exports/{export_file}/download")
def download_export(
    export_file: str,
    filename: str = Query(...),
    session: Session = Depends(get_session),
    user_id: int = Depends(get_current_user_id),
) -> FileResponse:
    try:
        path = ExportService(session).resolve_export_path(export_file, user_id=user_id)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail="导出文件不存在") from exc
    return FileResponse(
        path=path,
        media_type=_media_type(path),
        filename=Path(unquote(filename)).name,
    )


def _media_type(path: Path) -> str:
    if path.suffix == ".json":
        return "application/json"
    if path.suffix == ".md":
        return "text/markdown"
    if path.suffix == ".zip":
        return "application/zip"
    return "application/octet-stream"


@router.post("/api/projects/{project_id}/export-pptx", response_model=ExportResponse)
async def export_pptx(
    project_id: str,
    session: Session = Depends(get_session),
    user_id: int = Depends(get_current_user_id),
) -> ExportResponse:
    from app.config import settings
    from app.models.project import ProjectRecord

    project_service = ProjectService(session)
    data = project_service.get_project_data(project_id, user_id=user_id)

    slides_with_images = [s for s in data.slides if s.image_url]
    if not slides_with_images:
        raise HTTPException(status_code=400, detail="没有已生成的图片，无法导出 PPT")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    failed_slides: list[int] = []
    async with httpx.AsyncClient(timeout=120.0, follow_redirects=True) as client:
        for slide in sorted(slides_with_images, key=lambda s: s.slide_no):
            try:
                image_url = slide.image_url
                if image_url.startswith("data:"):
                    header, b64data = image_url.split(",", 1)
                    image_bytes = base64.b64decode(b64data)
                elif image_url.startswith(IMAGES_URL_PREFIX):
                    # 同进程内的落盘图：跳过 httpx 自调用，直接读盘。
                    local = resolve_local_path(image_url)
                    if not local or not local.exists():
                        raise FileNotFoundError(image_url)
                    image_bytes = local.read_bytes()
                else:
                    resp = await client.get(image_url)
                    resp.raise_for_status()
                    image_bytes = resp.content
                image_stream = BytesIO(image_bytes)
                pptx_slide = prs.slides.add_slide(blank_layout)
                pptx_slide.shapes.add_picture(image_stream, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
                if slide.speech_script and slide.speech_script.strip():
                    pptx_slide.notes_slide.notes_text_frame.text = slide.speech_script.strip()
            except Exception as exc:
                logger.warning("[export-pptx] slide {} image failed: {}", slide.slide_no, exc)
                failed_slides.append(slide.slide_no)
                continue

    if not prs.slides:
        detail = f"所有图片处理失败，无法生成 PPT。失败页码：{failed_slides}"
        raise HTTPException(status_code=502, detail=detail)

    export_dir = settings.storage_dir / "exports"
    export_dir.mkdir(parents=True, exist_ok=True)
    export_id = __import__('uuid').uuid4().hex[:12]
    path = export_dir / f"{project_id}__{export_id}.pptx"

    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    path.write_bytes(pptx_bytes.getvalue())

    project_record = session.get(ProjectRecord, project_id)
    project_title = project_record.title if project_record else data.project_id
    safe_title = re.sub(r"[\\/:*?\"<>|\s]+", "-", project_title.strip())
    safe_title = re.sub(r"-+", "-", safe_title).strip("-.")
    safe_title = safe_title[:80] or "material"
    filename = f"{safe_title}-slides.pptx"

    return ExportResponse(
        filename=filename,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        download_url=f"/api/exports/{path.name}/download?filename={filename}",
    )
